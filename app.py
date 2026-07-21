import streamlit as st
import os
import smtplib
import io
import speech_recognition as sr
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from PIL import Image
from pptx import Presentation

# --- SPEECH RECOGNITION ENGINE ---
def transcribe_actual_audio(audio_bytes, target_language_code):
    """Processes real microphone audio bytes and extracts spoken words using the selected regional language code."""
    recognizer = sr.Recognizer()
    
    try:
        audio_file = io.BytesIO(audio_bytes)
        with sr.AudioFile(audio_file) as source:
            recognizer.adjust_for_ambient_noise(source, duration=0.5)
            audio_data = recognizer.record(source)
            
            # Pass selected regional language code (e.g. 'hi-IN', 'te-IN') to Google Speech engine
            actual_text = recognizer.recognize_google(audio_data, language=target_language_code)
            return f"Grievance: {actual_text}"
            
    except sr.UnknownValueError:
        return "Grievance Status: Audio captured, but the AI could not clearly identify the spoken words. Please speak closer to your microphone."
    except sr.RequestError as e:
        return f"Speech Core Interruption: Network connection error ({e}). Please check your internet connection."
    except Exception as e:
        return f"Processing Error: Could not parse audio format. Ensure your recording contains clear voice input."

# --- SYSTEM UI CONFIGURATION ---
st.set_page_config(page_title="Grameena Seva App", page_icon="🌾", layout="centered")

st.markdown("""
    <style>
    .big-button { font-size:24px !important; font-weight: bold; }
    .stButton>button { width: 100%; height: 60px; background-color: #E2725B; color: white; border-radius: 10px; }
    </style>
    """, unsafe_allow_html=True)

# --- HEADER SECTION ---
st.title("🌾 Grameena Seva AI Hub")
st.write("Bridging the Linguistic Gap for Rural India")
st.write("---")

# Language code mapping for Indian regional languages
languages_map = {
    "English": "en-IN",
    "Hindi (हिन्दी)": "hi-IN",
    "Telugu (తెలుగు)": "te-IN",
    "Tamil (தமிழ்)": "ta-IN",
    "Kannada (ಕನ್ನಡ)": "kn-IN",
    "Marathi (मराठी)": "mr-IN",
    "Bengali (বাংলা)": "bn-IN",
    "Gujarati (ગુજરાતી)": "gu-IN",
    "Malayalam (മലയാളം)": "ml-IN",
    "Punjabi (ਪੰਜਾਬੀ)": "pa-IN",
    "Urdu (اُردُو)": "ur-IN"
}

selected_ui_lang = st.selectbox("🌐 Choose your language / भाषा चुनें", list(languages_map.keys()))
lang_code = languages_map[selected_ui_lang]

st.write(f"App set to: **{selected_ui_lang}** (Code: `{lang_code}`)")
st.write("---")

tab1, tab2 = st.tabs(["🎙️ Talk to Mitra (Voice)", "📷 Scan Documents (OCR)"])

# --- SMTP EMAIL DISPATCH SYSTEM ---
def dispatch_grievance_email(original_lang, transcribed_text, citizen_name, citizen_address):
    """Securely transmits the grievance and citizen profile via SMTP using Streamlit Secrets."""
    try:
        sender_email = st.secrets["SYSTEM_ALERT_EMAIL"]
        sender_password = st.secrets["SYSTEM_ALERT_PASSWORD"]
        receiver_email = st.secrets["GRIEVANCE_OFFICER_EMAIL"]
    except Exception as e:
        st.sidebar.error(f"Secrets missing: {e}")
        return False
    
    msg = MIMEMultipart()
    msg['From'] = sender_email
    msg['To'] = receiver_email
    msg['Subject'] = f"URGENT: New Grievance Filed by {citizen_name} ({original_lang})"
    
    body = f"""
    Respected Officer,
    
    A new citizen grievance has been submitted via the Grameena Seva App platform.
    
    --- CITIZEN IDENTITY DETAILS ---
    Name of Citizen : {citizen_name}
    Village/Address : {citizen_address}
    
    --- FIELD REPORT ---
    Language Selected: {original_lang}
    {transcribed_text}
    
    --------------------------------------------------
    This is an automated operational dispatch. Please review and initiate resolution workflows.
    """
    msg.attach(MIMEText(body, 'plain'))
    
    try:
        server = smtplib.SMTP_SSL('smtp.gmail.com', 465)
        server.login(sender_email, sender_password)
        text = msg.as_string()
        server.sendmail(sender_email, receiver_email, text)
        server.quit()
        return True
    except Exception as e:
        st.sidebar.error(f"Mail Server Error: {e}")
        return False

# --- TAB 1: VOICE GRIEVANCE FORM ---
with tab1:
    st.subheader("📝 Citizen Information Form")
    farmer_name = st.text_input("👤 Enter Full Name / पूरा नाम दर्ज करें *")
    farmer_address = st.text_area("🏠 Enter Village & Address / गांव और पता दर्ज करें *", height=70)
    
    st.write("---")
    st.subheader("🎙️ Record Your Grievance")
    audio_file = st.audio_input("Press record and speak naturally in your chosen language")
    
    if audio_file:
        st.success("Audio captured successfully!")
        
        audio_bytes = audio_file.read()
        st.info(f"Processing voice input in {selected_ui_lang}...")
        
        # Transcribe audio using the selected language code
        translated_text = transcribe_actual_audio(audio_bytes, lang_code)
        
        st.subheader("📋 Output Verification:")
        st.text_area(label="", value=translated_text, height=120)
        
        # Validation Gate Check
        if farmer_name.strip() and farmer_address.strip():
            if st.button("🚀 Submit Grievance to Government Portal"):
                with st.spinner("Routing alert to the official department..."):
                    email_status = dispatch_grievance_email(
                        selected_ui_lang, 
                        translated_text, 
                        farmer_name, 
                        farmer_address
                    )
                    
                    if email_status:
                        st.success(f"🎉 Grievance filed for {farmer_name}! Sent to the government official.")
                    else:
                        st.error("Submission failed. Check system credentials.")
        else:
            st.warning("⚠️ Action Required: Please fill out both your Name and Address fields above to unlock the Submit button.")

# --- TAB 2: DOCUMENT OCR PLACEHOLDER ---
with tab2:
    st.subheader("Upload or Snap ID / Land Records")
    uploaded_image = st.camera_input("Take a photo of Document or Land Record")
    if uploaded_image:
        image = Image.open(uploaded_image)
        st.image(image, caption="Uploaded Document", use_container_width=True)
        st.info("Extracting data via Mobile OCR Engine...")

# --- POWERPOINT DECK GENERATOR ---
def build_pitch_deck():
    """Generates an in-memory PowerPoint presentation using python-pptx."""
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "🌾 Grameena Seva AI Hub"
    slide.placeholders[1].text = "Bridging the Linguistic Gap for Rural Indian Farmers\nCreated by Varun Sagar & Team"
    
    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "The Core Problem Statement"
    tf = slide.placeholders[1].text_frame
    tf.text = "• 65%+ of rural India relies heavily on regional native languages."
    tf.add_paragraph().text = "• Digital grievance channels remain heavily restricted to English-centric layouts."
    tf.add_paragraph().text = "• Keyboard barriers and administrative lag prevent simple problem resolution workflows."
    
    # Slide 3: Solution Framework
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Our Solution Framework"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Native Audio Layer: Farmers speak naturally into the regional interface."
    tf.add_paragraph().text = "• Speech-to-Text AI Pipeline: Dynamically converts acoustic data in real time."
    tf.add_paragraph().text = "• Verified Transcripts: Displays output for validation before dispatch."
    tf.add_paragraph().text = "• Automated SMTP Delivery: Securely emails reports directly to official inboxes."

    # Slide 4: Future Growth Focus
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Future Deployment Focus"
    tf = slide.placeholders[1].text_frame
    tf.text = "• OCR Integration: Real-time physical land record camera uploads."
    tf.add_paragraph().text = "• Bi-Directional Alerts: Re-translating officer replies back into native audio."
    tf.add_paragraph().text = "• Offline Gateway Fallback: Cellular SMS communications infrastructure."

    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# --- DOWNLOAD DASHBOARD AT BOTTOM ---
st.write("---")
st.subheader("📊 Presentation Exporter Dashboard")
st.write("Generate and download your high-impact project presentation deck below:")

presentation_data = build_pitch_deck()

st.download_button(
    label="📥 Download Pitch Deck (.pptx)",
    data=presentation_data,
    file_name="Grameena_Seva_Pitch_Deck.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    use_container_width=True
)
